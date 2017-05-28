import argparse, logging
from io import BytesIO
from time import time
from types import MethodType

import pyscreenshot
from pptx import Presentation
from pptx.util import Emu
from PIL import Image

def parse_args():
    parser = argparse.ArgumentParser(
        description=("Record a screencast as a"
                     " series of power-point slides")
    )

    parser.add_argument("file", metavar="file", type=str)
    parser.add_argument("-v", "--verbosity", action="store_true")
    return parser.parse_args()

def take_images(image_array):
    """
        append the images taken to an array
    """
    while True:
        image = pyscreenshot.grab()
        image_array.append(image)
        logging.info(
            "screenshot taken, image_array: [{}]".format(len(image_array))
        )

def resize_patch_image(image, size):
    """
        gets the image
        and it returns it resized and 'patched'
        to be processed by the pptx.add_picture()
    """
    image.thumbnail(size, Image.ANTIALIAS)
    # workaround because the add_picture expects an read()
    # to return the bytes
    b_image = BytesIO()
    image.save(b_image, format='PNG')
    b_image.read = MethodType(lambda self: self.getvalue(), b_image)
    return b_image

def add_slide_to_ppt(ppt, image):
    """
        append an slide to the presentation with
        the image
    """
    size = (
        Emu(ppt.slide_width).pt,
        Emu(ppt.slide_height).pt
    )
    blank_slide_layout = ppt.slide_layouts[6]
    left = top = Emu(0)
    slide = ppt.slides.add_slide(blank_slide_layout)

    image = resize_patch_image(image, size)

    logging.info(
        "adding image to presentation, slides: [{}]".format(len(ppt.slides)),
    )
    pic = slide.shapes.add_picture(
        image, left, top
    )

def build_ppt(image_array):
    """
        buidls the ppt from the image_queue
    """
    ppt = Presentation()
    for image in image_array:
        add_slide_to_ppt(ppt, image)
    return ppt

def record_screen():
    """
        take the screenshots, than returns the pttx with the
        slides
    """
    image_array = []
    time_start = time()
    try:
        print("recording...")
        take_images(image_array)
    except KeyboardInterrupt:
        print("finished recording")
        pass

    time_end = time()
    print("recorded {0:.3f} seconds".format(time_end - time_start))

    print("building power point...")
    return build_ppt(image_array)

if __name__ == "__main__":
    args = parse_args()
    logging.basicConfig(
        level=logging.INFO if args.verbosity else logging.ERROR
    )

    ppt = record_screen()

    print("saving to file")
    ppt.save(args.file)
