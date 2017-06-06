from unittest import TestCase
from copy import copy

from PIL import Image
from pptx import Presentation

import pptrecorder

class PptRecorderTest(TestCase):

    def setUp(self):
        self.image = Image.new('1', (100, 100))

    def test_resize_patch_image(self):
        image = self.image
        patched_image = pptrecorder.resize_patch_image(image, (50, 50))
        self.assertIsInstance(patched_image.read(), bytes)

    def test_add_slide_to_ppt(self):
        ppt = Presentation()
        image = self.image

        pptrecorder.add_slide_to_ppt(ppt, image)
        self.assertEqual(1, len(ppt.slides))

    def test_build_ppt(self):
        image_one = self.image
        image_two = copy(self.image)
        ppt = pptrecorder.build_ppt((image_one, image_two))
        self.assertEqual(2, len(ppt.slides))
